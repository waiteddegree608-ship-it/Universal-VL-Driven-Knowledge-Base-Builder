import os
import win32com.client
from pptx import Presentation
import fitz  # PyMuPDF
from openai import OpenAI
import time

# Placeholder for API Key - waiting for user input
API_KEY = "sk-glinbzmkuqfkveeqeiizatpkwseliggmflxybeqtmxrmcwhz" 
BASE_URL = "https://api.siliconflow.cn/v1"

def extract_text_from_pdf(filepath):
    text = ""
    try:
        doc = fitz.open(filepath)
        for page in doc:
            text += page.get_text()
    except Exception as e:
        print(f"Error reading PDF {filepath}: {e}")
    return text

def extract_text_from_pptx(filepath):
    text = ""
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX {filepath}: {e}")
    return text

def extract_text_from_ppt(filepath):
    # Use win32com to open PPT and extract text
    text = ""
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        # powerpoint.Visible = 1 # Optional
        presentation = powerpoint.Presentations.Open(filepath, WithWindow=False)
        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    if shape.TextFrame.HasText:
                        text += shape.TextFrame.TextRange.Text + "\n"
        presentation.Close()
        # powerpoint.Quit() # Don't quit if we want to reuse, but safer to quit at end of script
    except Exception as e:
        print(f"Error reading PPT {filepath}: {e}")
    return text

def extract_text_from_doc(filepath):
    text = ""
    try:
        word = win32com.client.Dispatch("Word.Application")
        doc = word.Documents.Open(filepath)
        text = doc.Content.Text
        doc.Close()
        # word.Quit()
    except Exception as e:
        print(f"Error reading DOC {filepath}: {e}")
    return text

def process_content(content):
    if not content.strip():
        return ""
    
    client = OpenAI(api_key=API_KEY, base_url=BASE_URL)
    
    # Simple chunking if too long (simplified)
    # Ideally should chunk intelligentlly
    chunks = [content[i:i+4000] for i in range(0, len(content), 4000)]
    
    knowledge_points = ""
    
    for i, chunk in enumerate(chunks):
        try:
            response = client.chat.completions.create(
                model="Qwen/Qwen2.5-7B-Instruct", # Example model, might need user to specify or default
                messages=[
                    {"role": "system", "content": "你是一个微观经济学专家。请从以下文本中提取详细的知识点，整理成Markdown格式。要求事无巨细，不要遗漏。"},
                    {"role": "user", "content": chunk}
                ]
            )
            data = response.choices[0].message.content
            knowledge_points += data + "\n\n"
            print(f"Processed chunk {i+1}/{len(chunks)}")
        except Exception as e:
            print(f"Error calling API: {e}")
            
    return knowledge_points

def main():
    base_dir = r"E:\workspace\ddl\微观经济学"
    output_file = r"E:\workspace\ddl\微观经济学\knowledge_base.md"
    
    ppt_dir = os.path.join(base_dir, "ppt")
    pdf_file = os.path.join(base_dir, "微观经济学-现代观点-范里安-练习册答案+英文版.pdf")
    
    files_to_process = []
    
    # Add PDF
    if os.path.exists(pdf_file):
        files_to_process.append(pdf_file)
        
    # Add PPTs
    if os.path.exists(ppt_dir):
        for f in os.listdir(ppt_dir):
            if f.endswith(('.ppt', '.pptx', '.doc')): # Included doc as seen in file list
                files_to_process.append(os.path.join(ppt_dir, f))
                
    full_knowledge_base = "# 微观经济学知识库\n\n"
    
    for file_path in files_to_process:
        print(f"Processing {file_path}...")
        file_content = ""
        if file_path.endswith(".pdf"):
            file_content = extract_text_from_pdf(file_path)
        elif file_path.endswith(".pptx"):
            file_content = extract_text_from_pptx(file_path)
        elif file_path.endswith(".ppt"):
            file_content = extract_text_from_ppt(file_path)
        elif file_path.endswith(".doc"):
            file_content = extract_text_from_doc(file_path)
            
        if file_content:
            print(f"Extracted {len(file_content)} chars from {os.path.basename(file_path)}")
            knowledge = process_content(file_content)
            full_knowledge_base += f"## Source: {os.path.basename(file_path)}\n\n{knowledge}\n---\n\n"
            
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(full_knowledge_base)
    
    print("Knowledge base created successfully.")

if __name__ == "__main__":
    main()
